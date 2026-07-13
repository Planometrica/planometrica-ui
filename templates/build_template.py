#!/usr/bin/env python3
"""Планометрика — брендовый .pptx-шаблон.

Палитра и шрифтовые пары берутся из дизайн-системы @planometrica/ui
(src/tokens/colors.ts, src/tokens/typography.ts). Канон шрифтов:
Unbounded — только нейминг «Планометрика», Manrope — заголовки,
JetBrains Mono — основной текст.

Тема (theme1.xml) патчится напрямую: цветовая схема + major/minor шрифты,
чтобы новые слайды наследовали бренд без ручной стилизации.
"""

import re
import shutil
import zipfile
from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

OUT = Path.home() / "Projects/Planometrica/planometrica-ui/templates/planometrica-brand.pptx"

# --- Палитра (src/tokens/colors.ts) ---
BRAND = {
    "primary": "0A4C76",
    "secondary": "1A7BB3",
    "green": "92CF93",
    "orange": "F59E0B",
}
LIGHT = {
    "background": "FFFFFF",
    "foreground": "0A4C76",
    "secondary": "F0F9FF",
    "mutedForeground": "64748B",
    "destructive": "EF4444",
    "border": "E2E8F0",
    "sidebar": "FAFAFA",
}
DARK = {
    "background": "0F172A",
    "card": "1E293B",
    "primary": "3B82F6",
    "secondary": "334155",
    "foreground": "F8FAFC",
    "mutedForeground": "94A3B8",
    "destructive": "7F1D1D",
}
CHART = ["0A4C76", "1A7BB3", "92CF93", "F59E0B", "EF4444"]

F_LOGO = "Unbounded"
F_HEAD = "Manrope"
F_BODY = "JetBrains Mono"

SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)


def c(hexs: str) -> RGBColor:
    return RGBColor.from_string(hexs)


def rect(slide, x, y, w, h, fill=None, line=None, shape=MSO_SHAPE.RECTANGLE, line_w=0.75):
    sp = slide.shapes.add_shape(shape, x, y, w, h)
    if fill:
        sp.fill.solid()
        sp.fill.fore_color.rgb = c(fill)
    else:
        sp.fill.background()
    if line:
        sp.line.color.rgb = c(line)
        sp.line.width = Pt(line_w)
    else:
        sp.line.fill.background()
    sp.shadow.inherit = False
    return sp


def text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space_after=None, line_spacing=None):
    """runs: список параграфов; параграф = список кортежей (text, font, size, color, bold)."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if space_after is not None:
            p.space_after = space_after
        if line_spacing is not None:
            p.line_spacing = line_spacing
        for (t, fname, size, color, bold) in para:
            r = p.add_run()
            r.text = t
            r.font.name = fname
            r.font.size = Pt(size)
            r.font.color.rgb = c(color)
            r.font.bold = bold
    return tb


def wordmark(slide, x, y, color, size=14):
    text(slide, x, y, Inches(4), Inches(0.4),
         [[("Планометрика", F_LOGO, size, color, False)]])


def footer(slide, page):
    wordmark(slide, Inches(0.6), Inches(7.02), BRAND["primary"], size=9)
    text(slide, Inches(12.2), Inches(7.02), Inches(0.6), Inches(0.3),
         [[(f"{page:02d}", F_BODY, 10, LIGHT["mutedForeground"], False)]],
         align=PP_ALIGN.RIGHT)


def swatch(slide, x, y, w, h, hexs, name, extra=None, dark_text=False, border=None):
    rect(slide, x, y, w, h, fill=hexs, line=border)
    label_color = DARK["foreground"] if dark_text else "0F172A"
    lines = [[(name, F_HEAD, 11, label_color, True)],
             [(f"#{hexs}", F_BODY, 10, LIGHT["mutedForeground"], False)]]
    if extra:
        lines.append([(extra, F_BODY, 9, LIGHT["mutedForeground"], False)])
    text(slide, x, y + h + Inches(0.06), w, Inches(0.7), lines)


prs = Presentation()
prs.slide_width, prs.slide_height = SLIDE_W, SLIDE_H
blank = prs.slide_layouts[6]


def bg(slide, hexs):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = c(hexs)


# ============ 1. Титульный (тёмный, brand primary) ============
s = prs.slides.add_slide(blank)
bg(s, BRAND["primary"])
wordmark(s, Inches(0.6), Inches(0.55), "FFFFFF", size=16)
rect(s, Inches(0.6), Inches(3.0), Inches(0.9), Inches(0.07), fill=BRAND["green"])
text(s, Inches(0.6), Inches(3.3), Inches(11.5), Inches(1.6),
     [[("Заголовок презентации", F_HEAD, 44, "FFFFFF", True)]])
text(s, Inches(0.6), Inches(4.45), Inches(10), Inches(0.6),
     [[("Подзаголовок — краткое описание темы и контекста", F_HEAD, 18, LIGHT["secondary"], False)]])
text(s, Inches(0.6), Inches(6.6), Inches(8), Inches(0.4),
     [[("2026 · planometrica.pro", F_BODY, 12, BRAND["green"], False)]])
rect(s, Inches(12.13), Inches(0), Inches(1.2), Inches(0.22), fill=BRAND["orange"])
rect(s, Inches(12.73), Inches(0.22), Inches(0.6), Inches(0.22), fill=BRAND["green"])

# ============ 2. Разделитель секции (brand secondary) ============
s = prs.slides.add_slide(blank)
bg(s, BRAND["secondary"])
text(s, Inches(0.6), Inches(1.6), Inches(4), Inches(2.2),
     [[("01", F_HEAD, 120, "FFFFFF", True)]])
rect(s, Inches(0.72), Inches(4.15), Inches(0.9), Inches(0.07), fill=BRAND["orange"])
text(s, Inches(0.72), Inches(4.45), Inches(10.5), Inches(1),
     [[("Название раздела", F_HEAD, 36, "FFFFFF", True)]])
text(s, Inches(0.72), Inches(5.35), Inches(9), Inches(0.5),
     [[("Пояснение к разделу — основной текст в JetBrains Mono", F_BODY, 14, LIGHT["secondary"], False)]])

# ============ 3. Контентный слайд ============
s = prs.slides.add_slide(blank)
bg(s, LIGHT["background"])
text(s, Inches(0.6), Inches(0.55), Inches(11), Inches(0.7),
     [[("Заголовок раздела", F_HEAD, 28, BRAND["primary"], True)]])
rect(s, Inches(0.6), Inches(1.25), Inches(1.2), Inches(0.06), fill=BRAND["green"])
bullets = [
    "Основной текст набирается JetBrains Mono — канон дизайн-системы",
    "Заголовки и подзаголовки — Manrope",
    "Unbounded используется только для нейминга «Планометрика»",
    "Цвета — строго из палитры @planometrica/ui",
]
paras = [[(f"—  {b}", F_BODY, 14, DARK["background"], False)] for b in bullets]
text(s, Inches(0.6), Inches(1.75), Inches(11.9), Inches(3.5), paras,
     space_after=Pt(14), line_spacing=1.15)
text(s, Inches(0.6), Inches(6.2), Inches(11.9), Inches(0.4),
     [[("Примечание: подписи и сноски — JetBrains Mono, muted #64748B", F_BODY, 11,
        LIGHT["mutedForeground"], False)]])
footer(s, 3)

# ============ 4. Две карточки ============
s = prs.slides.add_slide(blank)
bg(s, LIGHT["background"])
text(s, Inches(0.6), Inches(0.55), Inches(11), Inches(0.7),
     [[("Сравнение в две колонки", F_HEAD, 28, BRAND["primary"], True)]])
rect(s, Inches(0.6), Inches(1.25), Inches(1.2), Inches(0.06), fill=BRAND["green"])
card_y, card_h = Inches(1.8), Inches(4.6)
for i, (title, accent) in enumerate([("Вариант A", BRAND["secondary"]),
                                     ("Вариант B", BRAND["orange"])]):
    x = Inches(0.6 + i * 6.3)
    rect(s, x, card_y, Inches(5.9), card_h, fill=LIGHT["secondary"],
         line=LIGHT["border"], shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s, x + Inches(0.4), card_y + Inches(0.4), Inches(0.5), Inches(0.07), fill=accent)
    text(s, x + Inches(0.4), card_y + Inches(0.65), Inches(5.1), Inches(0.5),
         [[(title, F_HEAD, 18, BRAND["primary"], True)]])
    text(s, x + Inches(0.4), card_y + Inches(1.3), Inches(5.1), Inches(2.8),
         [[("Текст карточки — JetBrains Mono 13pt. Фон карточки — "
            "semantic secondary #F0F9FF, рамка — border #E2E8F0.", F_BODY, 13,
            DARK["background"], False)]], line_spacing=1.2)
footer(s, 4)

# ============ 5. Палитра: бренд + графики ============
s = prs.slides.add_slide(blank)
bg(s, LIGHT["background"])
text(s, Inches(0.6), Inches(0.55), Inches(11), Inches(0.7),
     [[("Палитра — бренд и диаграммы", F_HEAD, 28, BRAND["primary"], True)]])
rect(s, Inches(0.6), Inches(1.25), Inches(1.2), Inches(0.06), fill=BRAND["green"])
names = [("primary", "Основной синий"), ("secondary", "Светлый синий"),
         ("green", "Зелёный (success)"), ("orange", "Оранжевый (warning)")]
for i, (key, label) in enumerate(names):
    swatch(s, Inches(0.6 + i * 3.15), Inches(1.8), Inches(2.85), Inches(1.7),
           BRAND[key], label, extra=f"brand.{key}")
text(s, Inches(0.6), Inches(4.6), Inches(11), Inches(0.4),
     [[("Диаграммы (chart.1–5)", F_HEAD, 16, BRAND["primary"], True)]])
for i, hexs in enumerate(CHART):
    swatch(s, Inches(0.6 + i * 2.5), Inches(5.15), Inches(2.2), Inches(0.9),
           hexs, f"chart.{i + 1}")
footer(s, 5)

# ============ 6. Палитра: семантика light/dark ============
s = prs.slides.add_slide(blank)
bg(s, LIGHT["background"])
text(s, Inches(0.6), Inches(0.55), Inches(11), Inches(0.7),
     [[("Палитра — семантические цвета", F_HEAD, 28, BRAND["primary"], True)]])
rect(s, Inches(0.6), Inches(1.25), Inches(1.2), Inches(0.06), fill=BRAND["green"])

light_items = [("background", "FFFFFF"), ("secondary / muted", "F0F9FF"),
               ("border / input", "E2E8F0"), ("sidebar", "FAFAFA"),
               ("mutedForeground", "64748B"), ("destructive", "EF4444"),
               ("foreground", "0A4C76")]
dark_items = [("background", "0F172A"), ("card / popover", "1E293B"),
              ("secondary / muted", "334155"), ("primary", "3B82F6"),
              ("mutedForeground", "94A3B8"), ("destructive", "7F1D1D"),
              ("foreground", "F8FAFC")]

text(s, Inches(0.6), Inches(1.6), Inches(5), Inches(0.4),
     [[("Light theme", F_HEAD, 16, BRAND["primary"], True)]])
text(s, Inches(6.9), Inches(1.6), Inches(5), Inches(0.4),
     [[("Dark theme", F_HEAD, 16, BRAND["primary"], True)]])
rect(s, Inches(6.75), Inches(2.1), Inches(6), Inches(4.9), fill=DARK["background"])
for i, (label, hexs) in enumerate(light_items):
    y = Inches(2.25 + i * 0.68)
    rect(s, Inches(0.6), y, Inches(0.5), Inches(0.5), fill=hexs,
         line=LIGHT["border"] if hexs in ("FFFFFF", "F0F9FF", "FAFAFA") else None)
    text(s, Inches(1.3), y + Inches(0.03), Inches(3.4), Inches(0.3),
         [[(label, F_HEAD, 12, "0F172A", True)]])
    text(s, Inches(4.9), y + Inches(0.05), Inches(1.4), Inches(0.3),
         [[(f"#{hexs}", F_BODY, 11, LIGHT["mutedForeground"], False)]])
for i, (label, hexs) in enumerate(dark_items):
    y = Inches(2.25 + i * 0.68)
    rect(s, Inches(7.0), y, Inches(0.5), Inches(0.5), fill=hexs,
         line=DARK["secondary"] if hexs == "0F172A" else None)
    text(s, Inches(7.7), y + Inches(0.03), Inches(3.2), Inches(0.3),
         [[(label, F_HEAD, 12, DARK["foreground"], True)]])
    text(s, Inches(11.1), y + Inches(0.05), Inches(1.4), Inches(0.3),
         [[(f"#{hexs}", F_BODY, 11, DARK["mutedForeground"], False)]])
footer(s, 6)

# ============ 7. Диаграмма ============
s = prs.slides.add_slide(blank)
bg(s, LIGHT["background"])
text(s, Inches(0.6), Inches(0.55), Inches(11), Inches(0.7),
     [[("Данные и диаграммы", F_HEAD, 28, BRAND["primary"], True)]])
rect(s, Inches(0.6), Inches(1.25), Inches(1.2), Inches(0.06), fill=BRAND["green"])
cd = CategoryChartData()
cd.categories = ["2022", "2023", "2024", "2025", "2026"]
cd.add_series("Показатель", (18, 32, 47, 65, 84))
gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED,
                        Inches(0.6), Inches(1.7), Inches(12.1), Inches(4.9), cd)
chart = gf.chart
chart.has_legend = False
plot = chart.plots[0]
plot.vary_by_categories = True
for i, point in enumerate(plot.series[0].points):
    point.format.fill.solid()
    point.format.fill.fore_color.rgb = c(CHART[i % len(CHART)])
for ax in (chart.category_axis, chart.value_axis):
    ax.tick_labels.font.name = F_BODY
    ax.tick_labels.font.size = Pt(11)
    ax.tick_labels.font.color.rgb = c(LIGHT["mutedForeground"])
footer(s, 7)

# ============ 8. Таблица ============
s = prs.slides.add_slide(blank)
bg(s, LIGHT["background"])
text(s, Inches(0.6), Inches(0.55), Inches(11), Inches(0.7),
     [[("Таблицы", F_HEAD, 28, BRAND["primary"], True)]])
rect(s, Inches(0.6), Inches(1.25), Inches(1.2), Inches(0.06), fill=BRAND["green"])
rows, cols = 5, 4
tbl = s.shapes.add_table(rows, cols, Inches(0.6), Inches(1.8),
                         Inches(12.1), Inches(3.6)).table
headers = ["Метрика", "2024", "2025", "2026"]
data = [["Выручка, млн ₽", "12,4", "18,9", "27,3"],
        ["Проекты", "86", "142", "210"],
        ["Партнёры", "14", "31", "58"],
        ["NPS", "62", "71", "78"]]
for j, h in enumerate(headers):
    cell = tbl.cell(0, j)
    cell.fill.solid()
    cell.fill.fore_color.rgb = c(BRAND["primary"])
    p = cell.text_frame.paragraphs[0]
    r = p.add_run(); r.text = h
    r.font.name = F_HEAD; r.font.size = Pt(13); r.font.bold = True
    r.font.color.rgb = c("FFFFFF")
for i, row in enumerate(data, start=1):
    for j, val in enumerate(row):
        cell = tbl.cell(i, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = c(LIGHT["secondary"] if i % 2 == 0 else "FFFFFF")
        p = cell.text_frame.paragraphs[0]
        r = p.add_run(); r.text = val
        r.font.name = F_BODY; r.font.size = Pt(12)
        r.font.color.rgb = c(DARK["background"])
footer(s, 8)

# ============ 9. Финальный (тёмный) ============
s = prs.slides.add_slide(blank)
bg(s, DARK["background"])
rect(s, Inches(5.97), Inches(2.55), Inches(1.4), Inches(0.07), fill=BRAND["green"])
text(s, Inches(1.7), Inches(2.9), Inches(10), Inches(0.9),
     [[("Планометрика", F_LOGO, 34, DARK["foreground"], False)]], align=PP_ALIGN.CENTER)
text(s, Inches(1.7), Inches(4.0), Inches(10), Inches(0.5),
     [[("Спасибо за внимание", F_HEAD, 20, DARK["mutedForeground"], False)]],
     align=PP_ALIGN.CENTER)
text(s, Inches(1.7), Inches(4.7), Inches(10), Inches(0.4),
     [[("planometrica.pro · hello@planometrica.pro", F_BODY, 13, BRAND["green"], False)]],
     align=PP_ALIGN.CENTER)

OUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(OUT)

# ============ Патч темы: палитра + шрифты ============
CLR_SCHEME = f'''<a:clrScheme name="Planometrica"><a:dk1><a:srgbClr val="{BRAND['primary']}"/></a:dk1><a:lt1><a:srgbClr val="FFFFFF"/></a:lt1><a:dk2><a:srgbClr val="{DARK['background']}"/></a:dk2><a:lt2><a:srgbClr val="F0F9FF"/></a:lt2><a:accent1><a:srgbClr val="{BRAND['primary']}"/></a:accent1><a:accent2><a:srgbClr val="{BRAND['secondary']}"/></a:accent2><a:accent3><a:srgbClr val="{BRAND['green']}"/></a:accent3><a:accent4><a:srgbClr val="{BRAND['orange']}"/></a:accent4><a:accent5><a:srgbClr val="{LIGHT['destructive']}"/></a:accent5><a:accent6><a:srgbClr val="{LIGHT['mutedForeground']}"/></a:accent6><a:hlink><a:srgbClr val="{BRAND['secondary']}"/></a:hlink><a:folHlink><a:srgbClr val="{LIGHT['mutedForeground']}"/></a:folHlink></a:clrScheme>'''

tmp = OUT.with_suffix(".tmp.pptx")
with zipfile.ZipFile(OUT) as zin, zipfile.ZipFile(tmp, "w", zipfile.ZIP_DEFLATED) as zout:
    for item in zin.infolist():
        data = zin.read(item.filename)
        if item.filename == "ppt/theme/theme1.xml":
            xml = data.decode("utf-8")
            xml = re.sub(r"<a:clrScheme .*?</a:clrScheme>", CLR_SCHEME, xml, flags=re.S)
            xml = re.sub(
                r'(<a:majorFont>\s*<a:latin typeface=")[^"]*(")',
                rf"\g<1>{F_HEAD}\g<2>", xml)
            xml = re.sub(
                r'(<a:minorFont>\s*<a:latin typeface=")[^"]*(")',
                rf"\g<1>{F_BODY}\g<2>", xml)
            xml = xml.replace('name="Office Theme"', 'name="Planometrica"')
            data = xml.encode("utf-8")
        zout.writestr(item, data)
shutil.move(tmp, OUT)

print(f"OK: {OUT}")
